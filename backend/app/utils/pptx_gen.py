from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
from typing import List, Dict, Any

def generate_pptx(slides_data: List[Dict[str, Any]]) -> io.BytesIO:
    """
    Generates a PowerPoint presentation with a dark-themed, modern design.
    """
    prs = Presentation()

    # Set 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Colors
    BG_COLOR = RGBColor(2, 6, 23)      # slate-950
    ACCENT_COLOR = RGBColor(37, 99, 235) # blue-600
    TEXT_MAIN = RGBColor(255, 255, 255) # white
    TEXT_DIM = RGBColor(203, 213, 225)  # slate-300
    CARD_BG = RGBColor(30, 41, 59)      # slate-800

    for i, slide_data in enumerate(slides_data):
        # Restore fallback for title
        title_str = slide_data.get("title") or f"Slide {i+1}"
        content_str = slide_data.get("content", "")
        layout_type = slide_data.get("layout_type", "text")

        # Use a blank layout for maximum control
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 1. Background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR

        # 2. Accent Bar at the bottom
        accent_bar = slide.shapes.add_shape(
            1, # Rectangle
            0, prs.slide_height - Inches(0.1), prs.slide_width, Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT_COLOR
        accent_bar.line.no_fill = True

        if i == 0:
            # Title Slide Layout
            # Main Title
            title_box = slide.shapes.add_textbox(
                0, prs.slide_height / 2 - Inches(1.2), prs.slide_width, Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            p = title_frame.paragraphs[0]
            run = p.add_run()
            run.text = title_str
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = TEXT_MAIN
            p.alignment = PP_ALIGN.CENTER

            # Subtitle / Content on Title Slide
            if content_str:
                subtitle_box = slide.shapes.add_textbox(
                    0, prs.slide_height / 2 + Inches(0.5), prs.slide_width, Inches(1)
                )
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.word_wrap = True
                p_sub = subtitle_frame.paragraphs[0]
                run_sub = p_sub.add_run()
                run_sub.text = content_str
                run_sub.font.size = Pt(24)
                run_sub.font.color.rgb = TEXT_DIM
                p_sub.alignment = PP_ALIGN.CENTER
        else:
            # Content Slide Layout
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            p = title_frame.paragraphs[0]
            run = p.add_run()
            run.text = title_str
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = TEXT_MAIN
            p.alignment = PP_ALIGN.LEFT

            # Title underline/accent
            title_line = slide.shapes.add_shape(
                1, # Rectangle
                Inches(0.5), Inches(0.9), Inches(2), Inches(0.03)
            )
            title_line.fill.solid()
            title_line.fill.fore_color.rgb = ACCENT_COLOR
            title_line.line.no_fill = True

            if layout_type == "architecture":
                # Render simple architecture diagram
                components = [line.strip() for line in content_str.split('\n') if line.strip()]
                if len(components) > 1:
                    box_width = Inches(2.2)
                    box_height = Inches(1.2)
                    start_x = Inches(0.5)
                    start_y = Inches(2.2)
                    spacing = Inches(0.3)

                    for j, comp in enumerate(components[:4]):
                        shape = slide.shapes.add_shape(
                            1, # Rectangle
                            start_x + (box_width + spacing) * j, start_y, box_width, box_height
                        )
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = ACCENT_COLOR
                        shape.line.color.rgb = TEXT_MAIN
                        shape.line.width = Pt(1.5)

                        tf = shape.text_frame
                        tf.word_wrap = True
                        p_comp = tf.paragraphs[0]
                        p_comp.text = comp
                        p_comp.font.size = Pt(14)
                        p_comp.font.bold = True
                        p_comp.font.color.rgb = TEXT_MAIN
                        p_comp.alignment = PP_ALIGN.CENTER
                else:
                    layout_type = "text"

            if layout_type == "grid":
                # Render cards
                items = [line.strip() for line in content_str.split('\n') if line.strip()]
                start_x = Inches(0.5)
                start_y = Inches(1.5)
                card_width = (prs.slide_width - Inches(1.5)) / 3
                card_height = Inches(1.6)

                for j, item in enumerate(items[:6]):
                    row = j // 3
                    col = j % 3
                    card = slide.shapes.add_shape(
                        1, # Rectangle
                        start_x + col * (card_width + Inches(0.2)),
                        start_y + row * (card_height + Inches(0.2)),
                        card_width, card_height
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = CARD_BG
                    card.line.color.rgb = ACCENT_COLOR
                    card.line.width = Pt(1.5)

                    tf = card.text_frame
                    tf.word_wrap = True
                    p_card = tf.paragraphs[0]
                    p_card.text = item.lstrip('-*• ').strip()
                    p_card.font.size = Pt(14)
                    p_card.font.color.rgb = TEXT_DIM
                    p_card.alignment = PP_ALIGN.LEFT

            elif layout_type == "text":
                # Content
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.2), prs.slide_width - Inches(1), prs.slide_height - Inches(1.8)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True

                lines = content_str.split('\n')
                for j, line in enumerate(lines):
                    line_text = line.strip()
                    if not line_text and j > 0: continue

                    if j == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()

                    if line_text.startswith(('-', '*', '•')):
                        p.level = 1
                        line_text = line_text[1:].strip()
                    else:
                        p.level = 0

                    run = p.add_run()
                    run.text = line_text
                    run.font.size = Pt(18)
                    run.font.color.rgb = TEXT_DIM
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(12)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io
