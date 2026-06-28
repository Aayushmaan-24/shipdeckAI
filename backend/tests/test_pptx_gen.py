import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from backend.app.utils.pptx_gen import generate_pptx

def test_generate_pptx_properties():
    slides_data = [
        {"slide": 1, "title": "Main Title", "content": "Welcome to ShipDeck Subtitle"},
        {"slide": 2, "title": "Features", "content": "- Feature 1\n- Feature 2"},
        {"slide": 3, "content": "Slide without title"}
    ]

    pptx_io = generate_pptx(slides_data)
    prs = Presentation(pptx_io)

    # Verify dimensions (16:9)
    assert prs.slide_width == Inches(10)
    assert prs.slide_height == Inches(5.625)

    # Verify slide count
    assert len(prs.slides) == 3

    # Colors to verify
    BG_COLOR_EXPECTED = RGBColor(2, 6, 23)
    ACCENT_COLOR_EXPECTED = RGBColor(37, 99, 235)
    TEXT_MAIN_EXPECTED = RGBColor(255, 255, 255)
    TEXT_DIM_EXPECTED = RGBColor(203, 213, 225)

    for i, slide in enumerate(prs.slides):
        # 1. Check background
        assert slide.background.fill.fore_color.rgb == BG_COLOR_EXPECTED

        # 2. Check accent bar (it's the first shape now)
        accent_bar = slide.shapes[0]
        assert accent_bar.fill.fore_color.rgb == ACCENT_COLOR_EXPECTED

        if i == 0:
            # Title slide
            title_box = slide.shapes[1]
            title_frame = title_box.text_frame
            assert title_frame.word_wrap is True

            p = title_frame.paragraphs[0]
            assert p.text == "Main Title"
            # Check run-level formatting
            run = p.runs[0]
            assert run.font.size == Pt(44)
            assert run.font.color.rgb == TEXT_MAIN_EXPECTED

            # Subtitle
            subtitle_box = slide.shapes[2]
            p_sub = subtitle_box.text_frame.paragraphs[0]
            assert p_sub.text == "Welcome to ShipDeck Subtitle"
            assert p_sub.runs[0].font.size == Pt(24)
        elif i == 1:
            # Content slide
            # Title
            title_box = slide.shapes[1]
            p_title = title_box.text_frame.paragraphs[0]
            assert p_title.text == "Features"
            assert p_title.runs[0].font.size == Pt(32)
            assert p_title.runs[0].font.color.rgb == TEXT_MAIN_EXPECTED

            # Title accent line (new)
            title_line = slide.shapes[2]
            assert title_line.fill.fore_color.rgb == ACCENT_COLOR_EXPECTED

            # Content
            content_box = slide.shapes[3]
            content_frame = content_box.text_frame
            assert content_frame.word_wrap is True
            assert len(content_frame.paragraphs) == 2
            assert content_frame.paragraphs[0].runs[0].font.size == Pt(18)
            assert content_frame.paragraphs[0].runs[0].font.color.rgb == TEXT_DIM_EXPECTED
            assert content_frame.paragraphs[0].level == 1
        elif i == 2:
            # Fallback title
            title_box = slide.shapes[1]
            assert title_box.text_frame.paragraphs[0].text == "Slide 3"

if __name__ == "__main__":
    test_generate_pptx_properties()
    print("PPTX Generation tests passed!")
